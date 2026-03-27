import re
from io import BytesIO
from typing import Dict, List

from docx import Document
from pptx import Presentation

from vision_engine import process_pdf

ALLOWED_EXTENSIONS = {".pdf", ".ppt", ".pptx", ".doc", ".docx", ".txt"}


def get_extension(filename: str) -> str:
    if not filename:
        return ""
    idx = filename.rfind(".")
    return filename[idx:].lower() if idx != -1 else ""


def extract_content_from_bytes(filename: str, raw: bytes) -> Dict[str, object]:
    """
    Normalized extraction output for all supported formats.
    Returns:
    {
      "mode": "text" | "image",
      "content": str | List[str]
    }
    """
    ext = get_extension(filename)
    if ext not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {ext or 'unknown'}")

    if ext == ".pdf":
        return process_pdf(BytesIO(raw))

    if ext == ".txt":
        return {"mode": "text", "content": _extract_txt(raw)}

    if ext == ".docx":
        return {"mode": "text", "content": _extract_docx(raw)}

    if ext == ".pptx":
        return {"mode": "text", "content": _extract_pptx(raw)}

    # Legacy Office formats (.doc/.ppt) - best effort fallback.
    return {"mode": "text", "content": _extract_legacy_binary_text(raw)}


def _extract_txt(raw: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            text = raw.decode(encoding)
            if text.strip():
                return text
        except Exception:
            continue
    raise ValueError("Unable to decode .txt file")


def _extract_docx(raw: bytes) -> str:
    doc = Document(BytesIO(raw))
    lines = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    if not lines:
        raise ValueError("No readable text found in .docx file")
    return "\n".join(lines)


def _extract_pptx(raw: bytes) -> str:
    presentation = Presentation(BytesIO(raw))
    chunks: List[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue

            text = str(getattr(text_frame, "text", "")).strip()
            if text:
                chunks.append(text)

    if not chunks:
        raise ValueError("No readable text found in .pptx file")
    return "\n".join(chunks)


def _extract_legacy_binary_text(raw: bytes) -> str:
    """Best-effort extraction for .doc/.ppt without external converters."""
    variants = []
    for encoding in ("utf-16le", "utf-8", "latin-1"):
        try:
            variants.append(raw.decode(encoding, errors="ignore"))
        except Exception:
            continue

    combined = "\n".join(variants)
    words = re.findall(r"[A-Za-z][A-Za-z0-9_\-]{2,}", combined)

    if not words:
        raise ValueError(
            "Unable to reliably extract text from legacy .doc/.ppt file. "
            "Please convert it to .docx/.pptx for best results."
        )

    # Keep extraction concise and readable.
    return " ".join(words[:6000])
